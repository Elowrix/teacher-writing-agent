from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class WorkbookInventory:
    file_name: str
    file_type: str
    sheet_name: str
    rows: int
    columns: int
    status: str
    note: str = ""


class ReportRunner:
    """Initial dashboard-free report production pipeline.

    This first implementation inventories the selected workbooks, records sheet
    dimensions, produces a quality-control workbook, and generates draft Word and
    PowerPoint outputs from the same inventory. Session-specific EPE parsers and
    the 12-question analytics are added on top of this stable entry point.
    """

    def __init__(self, project_root: Path, hmac_secret: str) -> None:
        self.project_root = project_root
        self.hmac_secret = hmac_secret

    def run(
        self,
        *,
        epe_files: Iterable[Path],
        registry_files: Iterable[Path],
        output_dir: Path,
    ) -> dict[str, Path]:
        output_dir.mkdir(parents=True, exist_ok=True)
        inventory = self._build_inventory(epe_files=epe_files, registry_files=registry_files)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        excel_path = output_dir / f"EPE_Veri_Kalite_Ozeti_{timestamp}.xlsx"
        word_path = output_dir / f"EPE_Yillar_Arasi_Analiz_Raporu_{timestamp}.docx"
        ppt_path = output_dir / f"EPE_Yonetim_Sunumu_{timestamp}.pptx"

        self._write_excel(inventory, excel_path)
        self._write_word(inventory, word_path)
        self._write_powerpoint(inventory, ppt_path)

        return {
            "Veri kalite özeti": excel_path,
            "Word raporu": word_path,
            "PowerPoint sunumu": ppt_path,
        }

    def _build_inventory(
        self,
        *,
        epe_files: Iterable[Path],
        registry_files: Iterable[Path],
    ) -> pd.DataFrame:
        rows: list[WorkbookInventory] = []
        for file_type, paths in (("EPE", epe_files), ("Öğrenci kütüğü", registry_files)):
            for path in paths:
                if not path.exists():
                    rows.append(
                        WorkbookInventory(
                            file_name=path.name,
                            file_type=file_type,
                            sheet_name="—",
                            rows=0,
                            columns=0,
                            status="HATA",
                            note="Dosya bulunamadı",
                        )
                    )
                    continue
                try:
                    workbook = pd.ExcelFile(path)
                except Exception as exc:  # noqa: BLE001
                    rows.append(
                        WorkbookInventory(
                            file_name=path.name,
                            file_type=file_type,
                            sheet_name="—",
                            rows=0,
                            columns=0,
                            status="HATA",
                            note=str(exc),
                        )
                    )
                    continue

                for sheet in workbook.sheet_names:
                    try:
                        frame = pd.read_excel(path, sheet_name=sheet, header=None)
                        rows.append(
                            WorkbookInventory(
                                file_name=path.name,
                                file_type=file_type,
                                sheet_name=str(sheet),
                                rows=int(frame.shape[0]),
                                columns=int(frame.shape[1]),
                                status="OK",
                            )
                        )
                    except Exception as exc:  # noqa: BLE001
                        rows.append(
                            WorkbookInventory(
                                file_name=path.name,
                                file_type=file_type,
                                sheet_name=str(sheet),
                                rows=0,
                                columns=0,
                                status="HATA",
                                note=str(exc),
                            )
                        )

        return pd.DataFrame([item.__dict__ for item in rows])

    @staticmethod
    def _write_excel(inventory: pd.DataFrame, path: Path) -> None:
        summary = (
            inventory.groupby(["file_type", "status"], dropna=False)
            .size()
            .reset_index(name="sheet_count")
        )
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            inventory.to_excel(writer, sheet_name="Kaynak Envanteri", index=False)
            summary.to_excel(writer, sheet_name="Özet", index=False)

    @staticmethod
    def _write_word(inventory: pd.DataFrame, path: Path) -> None:
        document = Document()
        title = document.add_heading("EPE Yıllar Arası Analiz Raporu", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        paragraph = document.add_paragraph(
            "Bu sürüm, seçilen EPE sonuç dosyaları ve öğrenci kütüklerinin kaynak "
            "envanterini doğrular. Oturum bazlı analizler ve 12 temel sorunun sonuçları "
            "aynı rapor üretim hattına eklenecektir."
        )
        paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        document.add_heading("1. Kaynak Dosya Özeti", level=1)
        file_summary = (
            inventory.groupby(["file_type", "file_name"], dropna=False)
            .agg(sheet_count=("sheet_name", "count"), error_count=("status", lambda s: int((s == "HATA").sum())))
            .reset_index()
        )
        table = document.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        headers = ["Kaynak türü", "Dosya", "Sayfa sayısı", "Hatalı sayfa"]
        for index, text in enumerate(headers):
            table.rows[0].cells[index].text = text
        for _, row in file_summary.iterrows():
            cells = table.add_row().cells
            cells[0].text = str(row["file_type"])
            cells[1].text = str(row["file_name"])
            cells[2].text = str(row["sheet_count"])
            cells[3].text = str(row["error_count"])

        document.add_heading("2. Sonraki Analiz Katmanı", level=1)
        for item in (
            "PASS/FAIL oranlarının aynı EPE yuvası içinde yıllar arasında karşılaştırılması",
            "Near-miss ve sınırda geçen öğrencilerin kendi karar eşiklerine göre belirlenmesi",
            "ELL–ELT ve diğer lisans programlarının ayrı raporlanması",
            "Fakülte ve burs kırılımları",
            "Booklet, Writing, Speaking ve Productive beceri profilleri",
            "Aynı doğrulanmış tablolardan Word ve PowerPoint üretimi",
        ):
            document.add_paragraph(item, style="List Bullet")

        document.save(path)

    @staticmethod
    def _write_powerpoint(inventory: pd.DataFrame, path: Path) -> None:
        presentation = Presentation()

        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "EPE Yıllar Arası Analiz"
        slide.placeholders[1].text = "Kaynak doğrulama ve rapor üretim aracı"

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Yüklenen Kaynaklar"
        summary = (
            inventory.groupby("file_type", dropna=False)["file_name"]
            .nunique()
            .reset_index(name="file_count")
        )
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(4.5))
        frame = textbox.text_frame
        frame.clear()
        for idx, row in summary.iterrows():
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = f"{row['file_type']}: {row['file_count']} dosya"
            paragraph.font.size = Pt(24)
            paragraph.space_after = Pt(12)

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Raporlama Omurgası"
        textbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(8.3), Inches(4.8))
        frame = textbox.text_frame
        items = [
            "Kaynak dosya ve sayfa doğrulaması",
            "Oturum bazlı harmonizasyon",
            "PASS/FAIL, eşik bandı ve near-miss analizleri",
            "Fakülte, burs ve beceri profilleri",
            "Word raporu ve yönetim sunumu",
        ]
        for idx, item in enumerate(items):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.font.size = Pt(22)
            paragraph.level = 0

        presentation.save(path)
