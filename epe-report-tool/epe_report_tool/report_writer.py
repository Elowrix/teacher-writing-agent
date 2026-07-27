from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory

import matplotlib.pyplot as plt
import pandas as pd
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


RATE_FIELDS = {
    "PASS_rate", "Wilson_lower", "Wilson_upper", "near_miss_rate_among_FAIL",
    "borderline_rate_among_PASS", "PASS_rate_among_finalized", "absent_rate_among_listed",
    "unresolved_rate_among_listed", "registry_match_rate", "within_5_rate", "within_10_rate",
    "FAIL_0_5_rate", "FAIL_0_10_rate", "PASS_0_5_rate", "PASS_0_10_rate",
}
PERCENT_POINT_FIELDS = {
    "receptive_pct", "productive_pct", "writing_pct", "speaking_pct",
    "productive_minus_receptive", "writing_minus_speaking",
}
INTEGER_HINTS = ("_N", " N", "PASS", "FAIL", "rows", "columns", "row_n", "_n")


QUESTIONS = (
    ("Q1", "PASS/FAIL oranları yıllar içinde nasıl değişti?"),
    ("Q2", "Kalanların ne kadarı eşiğin ilk 5 puanı içinde?"),
    ("Q3", "Geçenlerin ne kadarı sınırda geçti?"),
    ("Q4", "ELL–ELT ile diğer lisans programlarının sonuçları nasıl farklılaşıyor?"),
    ("Q5", "Near-miss hangi fakültelerde yoğunlaşıyor?"),
    ("Q6–Q7", "Burs gruplarının sonuçları fakültelere göre nasıl değişiyor?"),
    ("Q8–Q10", "Eşik çevresindeki öğrencilerin beceri profilleri nasıl farklılaşıyor?"),
    ("Q11", "Giriş ve dönem içi EPE’lerde toplam puan ve beceri profilleri nasıl farklılaşıyor?"),
    ("Q12", "Öğrenciler eşik çevresinde mi yoğunlaşıyor?"),
)


def _tr_decimal(value: float, decimals: int = 1) -> str:
    return f"{float(value):.{decimals}f}".replace(".", ",")


def pct(value: object, *, suppress: bool = False) -> str:
    if suppress or value is None or pd.isna(value):
        return "—"
    return f"%{_tr_decimal(float(value) * 100, 1)}"


def number(value: object, decimals: int = 1) -> str:
    if value is None or pd.isna(value):
        return "—"
    try:
        return _tr_decimal(float(value), decimals)
    except (TypeError, ValueError):
        return str(value)


def integer(value: object) -> str:
    if value is None or pd.isna(value):
        return "—"
    try:
        return str(int(round(float(value))))
    except (TypeError, ValueError):
        return str(value)


def scholarship(value: object) -> str:
    if value is None or pd.isna(value):
        return "—"
    return f"%{integer(value)}"


def _small_cell_suppressed(row: pd.Series, field: str) -> bool:
    flag = str(row.get("cell_size_flag", ""))
    if not flag.startswith("N<5"):
        return False
    return field in RATE_FIELDS or field.endswith("_rate") or field == "PASS_rate"


def format_value(field: str, value: object, row: pd.Series | None = None) -> str:
    suppress = row is not None and _small_cell_suppressed(row, field)
    if suppress:
        return "—"
    if field in RATE_FIELDS or field.endswith("_rate") or field.startswith("Wilson"):
        return pct(value)
    if field == "scholarship":
        return scholarship(value)
    if field in PERCENT_POINT_FIELDS or field.endswith("_mean") or field.endswith("_median"):
        return number(value, 1)
    if isinstance(value, bool):
        return "Evet" if value else "Hayır"
    if field == "N" or field.endswith("_N") or field in {"PASS", "FAIL", "ABSENT"}:
        return integer(value)
    if isinstance(value, float):
        return number(value, 1)
    return "—" if value is None or pd.isna(value) else str(value)


def _set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def _set_repeat_table_header(row) -> None:
    tr_pr = row._tr.get_or_add_trPr()
    tbl_header = OxmlElement("w:tblHeader")
    tbl_header.set(qn("w:val"), "true")
    tr_pr.append(tbl_header)


def _add_docx_table(
    document: Document,
    frame: pd.DataFrame,
    columns: list[tuple[str, str]],
    max_rows: int | None = None,
    landscape: bool = False,
) -> None:
    if frame.empty:
        document.add_paragraph("Bu bölüm için kullanılabilir veri bulunmamaktadır.")
        return
    shown = frame if max_rows is None else frame.head(max_rows)
    table = document.add_table(rows=1, cols=len(columns))
    table.style = "Table Grid"
    table.autofit = True
    header = table.rows[0]
    _set_repeat_table_header(header)
    for idx, (_, label) in enumerate(columns):
        cell = header.cells[idx]
        cell.text = label
        _set_cell_shading(cell, "1F4E78")
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for paragraph in cell.paragraphs:
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = Pt(8.5)
    for _, row in shown.iterrows():
        cells = table.add_row().cells
        for idx, (field, _) in enumerate(columns):
            cells[idx].text = format_value(field, row.get(field), row)
            cells[idx].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for paragraph in cells[idx].paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if idx > 0 else WD_ALIGN_PARAGRAPH.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(8)
    if max_rows is not None and len(frame) > max_rows:
        document.add_paragraph(
            f"Tablonun ilk {max_rows} satırı gösterilmiştir. Tam tablo Excel teknik ekindedir."
        )


def _configure_document(document: Document) -> None:
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.65)
    section.right_margin = Inches(0.65)
    styles = document.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(9.5)
    for style_name in ("Title", "Heading 1", "Heading 2"):
        styles[style_name].font.name = "Arial"
        styles[style_name].font.color.rgb = RGBColor(31, 78, 121)


def _add_finding(document: Document, text: str) -> None:
    p = document.add_paragraph()
    p.paragraph_format.space_after = Pt(5)
    run = p.add_run("Temel bulgu: ")
    run.bold = True
    run.font.color.rgb = RGBColor(31, 78, 121)
    p.add_run(text)


def _result_finding(frame: pd.DataFrame) -> str:
    if frame.empty:
        return "Kullanılabilir sonuç bulunmamaktadır."
    best = frame.loc[frame["PASS_rate"].idxmax()]
    worst = frame.loc[frame["PASS_rate"].idxmin()]
    return (
        f"En yüksek PASS oranı {best['analysis_exam_id']} oturumunda {pct(best['PASS_rate'])}; "
        f"en düşük oran {worst['analysis_exam_id']} oturumunda {pct(worst['PASS_rate'])} olarak görülmektedir."
    )


def _near_finding(frame: pd.DataFrame, limitations: pd.DataFrame) -> str:
    if frame.empty:
        return "Near-miss sonucu bulunmamaktadır."
    usable = frame[frame["FAIL_N"].gt(0)]
    if usable.empty:
        text = "Finalized FAIL kaydı bulunan bir oturum yoktur."
    else:
        peak = usable.loc[usable["near_miss_rate_among_FAIL"].fillna(-1).idxmax()]
        text = (
            f"FAIL içinde en yüksek near-miss oranı {peak['analysis_exam_id']} oturumunda "
            f"{pct(peak['near_miss_rate_among_FAIL'])} düzeyindedir."
        )
    if not limitations.empty:
        exams = ", ".join(limitations["analysis_exam_id"].astype(str))
        text += f" {exams} için borderline süreç sınırlılığı nedeniyle sıfır değerleri doğal dağılım olarak yorumlanmamalıdır."
    return text


def _threshold_group_finding(frame: pd.DataFrame) -> str:
    if frame.empty:
        return "Eşik grubu karşılaştırması bulunmamaktadır."
    pivot = frame.pivot_table(
        index="analysis_exam_id", columns="threshold_group", values="PASS_rate", aggfunc="first"
    )
    diffs = []
    for exam_id, row in pivot.iterrows():
        if {"ELL–ELT", "Other undergraduate"}.issubset(row.dropna().index):
            diffs.append((exam_id, (row["ELL–ELT"] - row["Other undergraduate"]) * 100))
    if not diffs:
        return "İki eşik grubunun birlikte bulunduğu karşılaştırılabilir oturum yoktur."
    exam_id, diff = max(diffs, key=lambda item: abs(item[1]))
    direction = "ELL–ELT lehine" if diff > 0 else "diğer lisans programları lehine"
    return f"En büyük grup farkı {exam_id} oturumunda {abs(diff):.1f} puan ve {direction} görünmektedir.".replace(".", ",")


def _faculty_near_finding(frame: pd.DataFrame) -> str:
    if frame.empty:
        return "Fakülte near-miss dağılımı üretilememiştir."
    comparable = frame[(frame["FAIL_N"] >= 10) & frame["near_miss_rate_among_FAIL"].notna()]
    if comparable.empty:
        return "FAIL N≥10 koşulunu sağlayan fakülte hücresi bulunmadığı için karşılaştırma yapılmamıştır."
    peak = comparable.loc[comparable["near_miss_rate_among_FAIL"].idxmax()]
    unknown = int(frame.loc[frame["faculty"].eq("Unknown"), "N"].sum())
    note = (
        f"Karşılaştırılabilir hücrelerde en yüksek oran {peak['analysis_exam_id']} / {peak['faculty']} grubunda "
        f"{pct(peak['near_miss_rate_among_FAIL'])} düzeyindedir."
    )
    if unknown:
        note += f" Fakültesi bilinmeyen toplam {unknown} kayıt nedeniyle sonuçlar kapsam notuyla okunmalıdır."
    return note


def _skill_finding(frame: pd.DataFrame) -> str:
    if frame.empty:
        return "Tam beceri verisi bulunan eşik çevresi grubu yoktur."
    comparable = frame[frame["N"] >= 10]
    if comparable.empty:
        return "N≥10 koşulunu sağlayan beceri profili yoktur."
    strongest = comparable.loc[comparable["productive_minus_receptive"].idxmax()]
    weakest = comparable.loc[comparable["productive_minus_receptive"].idxmin()]
    return (
        f"Productive–Receptive farkı en yüksek {strongest['analysis_exam_id']} / {strongest['threshold_side']} "
        f"grubunda {number(strongest['productive_minus_receptive'])} puan; en düşük "
        f"{weakest['analysis_exam_id']} / {weakest['threshold_side']} grubunda "
        f"{number(weakest['productive_minus_receptive'])} puandır."
    )


def _exam_type_finding(frame: pd.DataFrame) -> str:
    if frame.empty:
        return "Giriş/dönem içi karşılaştırması üretilememiştir."
    years = frame["academic_year"].nunique()
    return (
        f"Giriş ve dönem içi EPE profilleri {years} akademik yıl için ayrı N, EPE ortalama/medyanı ve beceri "
        "ortalamalarıyla gösterilmiştir; farklı öğrenci popülasyonları nedeniyle betimsel okunmalıdır."
    )


def _concentration_finding(frame: pd.DataFrame) -> str:
    if frame.empty:
        return "Eşik yoğunluğu hesabı üretilememiştir."
    peak = frame.loc[frame["within_5_rate"].idxmax()]
    return (
        f"Eşiğin ±5 puanında en yüksek yoğunluk {peak['analysis_exam_id']} oturumunda "
        f"{pct(peak['within_5_rate'])} olarak görülmektedir."
    )


def _create_pass_chart(result_table: pd.DataFrame, path: Path) -> None:
    chart = result_table.copy()
    fig, ax = plt.subplots(figsize=(10, 5.2))
    values = chart["PASS_rate"] * 100
    ax.bar(chart["analysis_exam_id"].astype(str), values)
    ax.set_ylabel("PASS (%)")
    ax.set_ylim(0, 100)
    ax.set_title("Oturum Bazında PASS Oranı")
    ax.tick_params(axis="x", rotation=45)
    for index, value in enumerate(values):
        ax.text(index, value + 1.2, f"%{value:.1f}".replace(".", ","), ha="center", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def _create_near_chart(near_table: pd.DataFrame, path: Path) -> None:
    chart = near_table.copy()
    fig, ax = plt.subplots(figsize=(10, 5.2))
    values = chart["near_miss_rate_among_FAIL"].fillna(0) * 100
    ax.bar(chart["analysis_exam_id"].astype(str), values)
    ax.set_ylabel("FAIL içinde near-miss (%)")
    ax.set_ylim(0, max(20, float(values.max()) + 8))
    ax.set_title("FAIL Grubu İçinde Eşiğin İlk 5 Puanında Kalanlar")
    ax.tick_params(axis="x", rotation=45)
    for index, value in enumerate(values):
        ax.text(index, value + 0.7, f"%{value:.1f}".replace(".", ","), ha="center", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def _create_faculty_near_chart(frame: pd.DataFrame, path: Path) -> None:
    chart = frame[(frame["FAIL_N"] >= 10) & frame["near_miss_rate_among_FAIL"].notna()].copy()
    fig, ax = plt.subplots(figsize=(10, 6))
    if chart.empty:
        ax.text(0.5, 0.5, "Karşılaştırılabilir fakülte hücresi yok", ha="center", va="center")
        ax.axis("off")
    else:
        chart["label"] = chart["analysis_exam_id"].astype(str) + " | " + chart["faculty"].astype(str)
        chart = chart.sort_values("near_miss_rate_among_FAIL").tail(15)
        ax.barh(chart["label"], chart["near_miss_rate_among_FAIL"] * 100)
        ax.set_xlabel("FAIL içinde near-miss (%)")
        ax.set_title("Karşılaştırılabilir Fakülte Hücrelerinde Near-miss")
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def _create_exam_type_chart(frame: pd.DataFrame, path: Path) -> None:
    chart = frame.pivot_table(index="academic_year", columns="exam_type", values="epe_total_mean", aggfunc="first")
    fig, ax = plt.subplots(figsize=(9, 5))
    chart.plot(kind="bar", ax=ax)
    ax.set_ylabel("Ortalama EPE Total")
    ax.set_title("Giriş ve Dönem İçi EPE Toplam Puan Ortalamaları")
    ax.tick_params(axis="x", rotation=0)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def _create_concentration_chart(frame: pd.DataFrame, path: Path) -> None:
    chart = frame.copy()
    fig, ax = plt.subplots(figsize=(10, 5.2))
    x = range(len(chart))
    ax.bar(x, chart["within_5_rate"] * 100, label="±5 puan")
    ax.plot(x, chart["within_10_rate"] * 100, marker="o", label="±10 puan")
    ax.set_xticks(list(x), chart["analysis_exam_id"].astype(str), rotation=45)
    ax.set_ylabel("Finalized kayıtlar içindeki oran (%)")
    ax.set_title("Eşik Çevresindeki Yoğunluk")
    ax.legend()
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def write_excel(path: Path, tables: dict[str, pd.DataFrame]) -> None:
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        for name, frame in tables.items():
            safe_name = name[:31]
            frame.to_excel(writer, sheet_name=safe_name, index=False)
            sheet = writer.book[safe_name]
            sheet.freeze_panes = "A2"
            sheet.auto_filter.ref = sheet.dimensions
            header_fill = PatternFill("solid", fgColor="1F4E78")
            header_font = Font(color="FFFFFF", bold=True)
            for cell in sheet[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            for column_index, column_name in enumerate(frame.columns, start=1):
                column_letter = get_column_letter(column_index)
                width = min(
                    max(
                        len(str(column_name)),
                        max((len(str(value)) for value in frame[column_name].head(200) if pd.notna(value)), default=0),
                    ) + 2,
                    45,
                )
                sheet.column_dimensions[column_letter].width = width
                if column_name in RATE_FIELDS or column_name.endswith("_rate") or column_name.startswith("Wilson"):
                    for cell in sheet[column_letter][1:]:
                        cell.number_format = "0.0%"
                elif column_name in PERCENT_POINT_FIELDS or column_name.endswith("_mean") or column_name.endswith("_median"):
                    for cell in sheet[column_letter][1:]:
                        cell.number_format = "0.0"
                elif column_name == "scholarship":
                    for cell in sheet[column_letter][1:]:
                        cell.number_format = '0"%"'
                elif column_name == "N" or column_name.endswith("_N") or column_name in {"PASS", "FAIL", "ABSENT"}:
                    for cell in sheet[column_letter][1:]:
                        cell.number_format = "0"
            sheet.sheet_view.showGridLines = False


def write_word(path: Path, tables: dict[str, pd.DataFrame], coverage_note: str) -> None:
    result_table = tables.get("Genel Sonuclar", pd.DataFrame())
    participation = tables.get("Katilim Ozeti", pd.DataFrame())
    near_table = tables.get("Near Miss", pd.DataFrame())
    limitations = tables.get("Near Miss Sinirlilik", pd.DataFrame())
    threshold_table = tables.get("Esik Gruplari", pd.DataFrame())
    faculty_near = tables.get("Fakulte Near Miss", pd.DataFrame())
    faculty_burs = tables.get("Fakulte Burs", pd.DataFrame())
    skill_table = tables.get("Beceri Profili", pd.DataFrame())
    exam_type = tables.get("Sinav Turu", pd.DataFrame())
    concentration = tables.get("Esik Yogunlugu", pd.DataFrame())
    exclusions = tables.get("Dislanan Kayitlar", pd.DataFrame())
    quality_issues = tables.get("Kalite Sorunlari", pd.DataFrame())
    validation_table = tables.get("Dogrulama", pd.DataFrame())
    coverage = tables.get("12 Soru Kapsami", pd.DataFrame())

    document = Document()
    _configure_document(document)
    title = document.add_heading("EPE Yıllar Arası Analiz Raporu", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = document.add_paragraph("Eylül, Ocak ve Haziran/Temmuz EPE Oturumları")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.add_heading("Yönetici Özeti", level=1)
    document.add_paragraph(
        "Bu rapor, aynı EPE yuvasındaki oturumları yıllar arasında karşılaştırır. Resmî PASS/FAIL sonucu "
        "esas alınmış; karar puanı öğrencinin kendi program eşiğine uzaklığı, bant ve near-miss analizi için kullanılmıştır."
    )
    document.add_paragraph(coverage_note)
    _add_finding(document, _result_finding(result_table))
    _add_finding(document, _near_finding(near_table, limitations))

    document.add_heading("Kapsam ve Yöntem", level=1)
    for text in (
        "Diğer lisans programları için karar eşiği 64,50; ELL–ELT grubu için 74,50'dir.",
        "Near-miss: resmî sonucu FAIL olan ve karar puanı kendi eşiğinin 0–5 puan altında bulunan öğrenci.",
        "Sınırda geçen: resmî sonucu PASS olan ve karar puanı kendi eşiğinin 0–5 puan üstünde bulunan öğrenci.",
        "Aynı öğrencinin aynı analiz oturumunda birden fazla kaydı varsa son sitting korunur.",
        "Beceri analizleri yalnız Booklet, Writing ve Speaking puanlarının üçü de bulunan kayıtlarda yapılır.",
        "N<5 hücrelerde oran bastırılır; N=5–9 dikkatli, N≥10 karşılaştırılabilir kabul edilir.",
        "Eylül oturumlarında yapısal boşluk bulunan near-miss bölgeleri ayrıca sınırlılık olarak işaretlenir.",
    ):
        document.add_paragraph(text, style="List Bullet")

    with TemporaryDirectory() as temp_dir:
        temp = Path(temp_dir)

        document.add_heading("Katılım ve Sonuç Kapsamı", level=1)
        document.add_paragraph(
            "Listed N, finalized sınava giren N, absent ve unresolved/other kayıtlar birbirinden ayrılmıştır. "
            "PASS oranı yalnız finalized PASS+FAIL paydasından hesaplanır."
        )
        _add_docx_table(document, participation, [
            ("analysis_exam_id", "Oturum"), ("listed_N", "Listed N"),
            ("absent_N", "Absent N"), ("finalized_exam_taker_N", "Finalized N"),
            ("PASS", "PASS"), ("FAIL", "FAIL"),
            ("unresolved_or_other_N", "Other/Unresolved"),
            ("PASS_rate_among_finalized", "PASS %"),
        ])

        chart_path = temp / "pass.png"
        if not result_table.empty:
            _create_pass_chart(result_table, chart_path)
            document.add_heading("Q1. PASS/FAIL Oranlarının Yıllar İçindeki Değişimi", level=1)
            document.add_picture(str(chart_path), width=Inches(7.0))
            _add_finding(document, _result_finding(result_table))
            _add_docx_table(document, result_table, [
                ("analysis_exam_id", "Oturum"), ("N", "N"), ("PASS", "PASS"),
                ("FAIL", "FAIL"), ("PASS_rate", "PASS oranı"),
                ("Wilson_lower", "%95 GA alt"), ("Wilson_upper", "%95 GA üst"),
            ])

        document.add_heading("Q2–Q3. Near-miss ve Sınırda Geçenler", level=1)
        if not near_table.empty:
            chart_path = temp / "near.png"
            _create_near_chart(near_table, chart_path)
            document.add_picture(str(chart_path), width=Inches(7.0))
        _add_finding(document, _near_finding(near_table, limitations))
        _add_docx_table(document, near_table, [
            ("analysis_exam_id", "Oturum"), ("FAIL_N", "FAIL N"),
            ("near_miss_N", "Near-miss N"), ("near_miss_rate_among_FAIL", "FAIL içinde %"),
            ("PASS_N", "PASS N"), ("borderline_pass_N", "Sınırda PASS N"),
            ("borderline_rate_among_PASS", "PASS içinde %"),
        ])
        if not limitations.empty:
            document.add_heading("Near-miss Yorum Sınırlılıkları", level=2)
            _add_docx_table(document, limitations, [
                ("analysis_exam_id", "Oturum"), ("severity", "Düzey"), ("note", "Açıklama"),
            ])

        document.add_heading("Q4. ELL–ELT ve Diğer Lisans Programları", level=1)
        _add_finding(document, _threshold_group_finding(threshold_table))
        _add_docx_table(document, threshold_table, [
            ("analysis_exam_id", "Oturum"), ("threshold_group", "Eşik grubu"),
            ("N", "N"), ("PASS", "PASS"), ("FAIL", "FAIL"), ("PASS_rate", "PASS oranı"),
        ])

        document.add_heading("Q5. Near-miss’in Fakültelere Göre Dağılımı", level=1)
        _add_finding(document, _faculty_near_finding(faculty_near))
        if not faculty_near.empty:
            chart_path = temp / "faculty_near.png"
            _create_faculty_near_chart(faculty_near, chart_path)
            document.add_picture(str(chart_path), width=Inches(7.0))
        _add_docx_table(document, faculty_near, [
            ("analysis_exam_id", "Oturum"), ("faculty", "Fakülte"),
            ("N", "Toplam N"), ("FAIL_N", "FAIL N"),
            ("near_miss_N", "Near-miss N"),
            ("near_miss_rate_among_FAIL", "FAIL içinde %"),
            ("cell_size_flag", "Hücre değerlendirmesi"),
        ])

        document.add_heading("Q6–Q7. Fakülte ve Burs Görünümü", level=1)
        document.add_paragraph(
            "N<5 hücrelerde oranlar bastırılmıştır. Bu bölüm özet görünüm sunar; tam fakülte×burs tablosu Excel teknik ekindedir."
        )
        _add_docx_table(document, faculty_burs, [
            ("analysis_exam_id", "Oturum"), ("faculty", "Fakülte"),
            ("scholarship", "Burs"), ("N", "N"), ("PASS_rate", "PASS oranı"),
            ("near_miss_N", "Near-miss N"),
            ("near_miss_rate_among_FAIL", "FAIL içinde near-miss"),
            ("cell_size_flag", "Hücre değerlendirmesi"),
        ], max_rows=60)

        document.add_heading("Q8–Q10. Eşik Çevresinde Beceri Profilleri", level=1)
        document.add_paragraph(
            "Productive puan Writing+Speaking; Receptive puan Booklet üzerinden yüzde puana dönüştürülmüştür. "
            "Ocak 2024 orijinal birleşik EPE dosyasıyla beceri analizine dahildir."
        )
        _add_finding(document, _skill_finding(skill_table))
        _add_docx_table(document, skill_table, [
            ("analysis_exam_id", "Oturum"), ("threshold_side", "Eşik tarafı"), ("N", "N"),
            ("receptive_pct", "Receptive %"), ("productive_pct", "Productive %"),
            ("writing_pct", "Writing %"), ("speaking_pct", "Speaking %"),
            ("productive_minus_receptive", "Prod.–Rec."),
            ("writing_minus_speaking", "Writing–Speaking"),
        ])

        document.add_heading("Q11. Giriş ve Dönem İçi EPE Profilleri", level=1)
        document.add_paragraph(
            "Giriş ve dönem içi sınavlar farklı öğrenci popülasyonlarını temsil eder; bu nedenle tablo nedensel başarı sıralaması değil, betimsel profil karşılaştırmasıdır."
        )
        _add_finding(document, _exam_type_finding(exam_type))
        if not exam_type.empty:
            chart_path = temp / "exam_type.png"
            _create_exam_type_chart(exam_type, chart_path)
            document.add_picture(str(chart_path), width=Inches(6.8))
        _add_docx_table(document, exam_type, [
            ("academic_year", "Akademik yıl"), ("exam_type", "Sınav türü"), ("N", "N"),
            ("epe_total_N", "EPE total N"), ("epe_total_mean", "EPE ort."),
            ("epe_total_median", "EPE medyan"), ("skill_N", "Beceri N"),
            ("receptive_pct", "Receptive %"), ("productive_pct", "Productive %"),
            ("writing_pct", "Writing %"), ("speaking_pct", "Speaking %"),
        ])

        document.add_heading("Q12. Eşik Çevresindeki Yoğunluk", level=1)
        _add_finding(document, _concentration_finding(concentration))
        if not concentration.empty:
            chart_path = temp / "concentration.png"
            _create_concentration_chart(concentration, chart_path)
            document.add_picture(str(chart_path), width=Inches(7.0))
        _add_docx_table(document, concentration, [
            ("analysis_exam_id", "Oturum"), ("scored_finalized_N", "Puanlı N"),
            ("within_5_N", "±5 N"), ("within_5_rate", "±5 %"),
            ("within_10_N", "±10 N"), ("within_10_rate", "±10 %"),
            ("FAIL_0_5_rate", "FAIL ilk 5 %"), ("FAIL_0_10_rate", "FAIL ilk 10 %"),
            ("PASS_0_5_rate", "PASS ilk 5 %"), ("PASS_0_10_rate", "PASS ilk 10 %"),
        ])

    document.add_heading("Dışlanan Kayıtların Özeti", level=1)
    document.add_paragraph(
        "Yüksek lisans, doktora ve dismissed kayıtları kaynaklardan silinmemiş; ana lisans analizinden çıkarılarak aşağıda toplu olarak gösterilmiştir."
    )
    _add_docx_table(document, exclusions, [
        ("analysis_exam_id", "Oturum"), ("exclusion_reason", "Dışlama nedeni"),
        ("N", "N"), ("PASS", "PASS"), ("FAIL", "FAIL"), ("ABSENT", "Absent"),
    ])

    document.add_heading("Veri Kalitesi ve Doğrulama", level=1)
    document.add_paragraph(
        "Aşağıdaki liste, tanınmayan kaynakları, beklenen N farklarını, finalized sonucu olmayan kayıtları ve kütük eşleşmesi sonrasında kalan Unknown fakülteleri açıkça gösterir."
    )
    _add_docx_table(document, quality_issues, [
        ("severity", "Düzey"), ("area", "Alan"), ("file_or_exam", "Dosya/Oturum"),
        ("issue", "Sorun"), ("impact", "Etkisi / kontrol notu"),
    ])
    document.add_heading("PASS/FAIL–Eşik Doğrulaması", level=2)
    _add_docx_table(document, validation_table, [
        ("analysis_exam_id", "Oturum"), ("result_threshold_validation", "Kontrol"), ("N", "N"),
    ])

    document.add_heading("12 Temel Sorunun Kapsam Kontrolü", level=1)
    _add_docx_table(document, coverage, [
        ("question_id", "Soru"), ("question", "Analiz sorusu"),
        ("source_table", "Kaynak tablo"), ("supported", "Yanıtlandı"),
    ])

    document.add_heading("Kurumsal Değerlendirme Alanı", level=1)
    document.add_paragraph(
        "Bu bölüm, sonuçların program hedefleri, sınav uygulamaları ve öğrenci destek kararları açısından kurum tarafından yorumlanması için ayrılmıştır."
    )
    document.add_paragraph("\n" * 4)
    document.save(path)


def _add_ppt_table(slide, frame: pd.DataFrame, columns: list[tuple[str, str]], max_rows: int = 8) -> None:
    if frame.empty:
        textbox = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(11.5), PptInches(3.5))
        textbox.text_frame.text = "Bu bölüm için kullanılabilir veri bulunmamaktadır."
        return
    shown = frame.head(max_rows)
    shape = slide.shapes.add_table(
        len(shown) + 1, len(columns), PptInches(0.35), PptInches(1.25), PptInches(12.6), PptInches(5.7)
    )
    table = shape.table
    for idx, (_, label) in enumerate(columns):
        table.cell(0, idx).text = label
    for row_idx, (_, row) in enumerate(shown.iterrows(), start=1):
        for col_idx, (field, _) in enumerate(columns):
            table.cell(row_idx, col_idx).text = format_value(field, row.get(field), row)
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_anchor = 1
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.size = PptPt(9 if row_idx else 10)
                paragraph.font.bold = row_idx == 0


def _new_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _new_content_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    return slide


def _add_ppt_finding(slide, text: str) -> None:
    box = slide.shapes.add_textbox(PptInches(0.7), PptInches(6.55), PptInches(12.0), PptInches(0.55))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = PptPt(15)
    tf.paragraphs[0].font.bold = True


def write_powerpoint(path: Path, tables: dict[str, pd.DataFrame], coverage_note: str) -> None:
    result_table = tables.get("Genel Sonuclar", pd.DataFrame())
    participation = tables.get("Katilim Ozeti", pd.DataFrame())
    near_table = tables.get("Near Miss", pd.DataFrame())
    limitations = tables.get("Near Miss Sinirlilik", pd.DataFrame())
    threshold_table = tables.get("Esik Gruplari", pd.DataFrame())
    faculty_near = tables.get("Fakulte Near Miss", pd.DataFrame())
    skill_table = tables.get("Beceri Profili", pd.DataFrame())
    exam_type = tables.get("Sinav Turu", pd.DataFrame())
    concentration = tables.get("Esik Yogunlugu", pd.DataFrame())
    quality_issues = tables.get("Kalite Sorunlari", pd.DataFrame())
    exclusions = tables.get("Dislanan Kayitlar", pd.DataFrame())

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    _new_title_slide(prs, "EPE Yıllar Arası Analiz", "12 temel soruya karar odaklı yanıt")

    slide = _new_content_slide(prs, "Kapsam ve Okuma İlkeleri")
    textbox = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.35), PptInches(11.8), PptInches(5.4))
    tf = textbox.text_frame
    principles = [
        coverage_note,
        "Karşılaştırmalar aynı EPE yuvası içinde yıllar arasında yapılır.",
        "N<5 oranları bastırılır; N=5–9 dikkatli okunur.",
        "Eylül near-miss sıfırları borderline süreç sınırlılığıyla birlikte yorumlanır.",
        "Resmî PASS/FAIL sonucu değiştirilmez.",
    ]
    for idx, item in enumerate(principles):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = PptPt(18)
        p.space_after = PptPt(10)

    slide = _new_content_slide(prs, "Katılım ve Sonuç Kapsamı")
    _add_ppt_table(slide, participation, [
        ("analysis_exam_id", "Oturum"), ("listed_N", "Listed"),
        ("absent_N", "Absent"), ("finalized_exam_taker_N", "Finalized"),
        ("PASS", "PASS"), ("FAIL", "FAIL"),
        ("unresolved_or_other_N", "Other"),
        ("PASS_rate_among_finalized", "PASS %"),
    ], max_rows=9)

    with TemporaryDirectory() as temp_dir:
        temp = Path(temp_dir)
        if not result_table.empty:
            chart = temp / "pass.png"
            _create_pass_chart(result_table, chart)
            slide = _new_content_slide(prs, "Q1 | PASS Oranlarının Yıllar İçindeki Değişimi")
            slide.shapes.add_picture(str(chart), PptInches(0.8), PptInches(1.2), width=PptInches(11.8))
            _add_ppt_finding(slide, _result_finding(result_table))

        if not near_table.empty:
            chart = temp / "near.png"
            _create_near_chart(near_table, chart)
            slide = _new_content_slide(prs, "Q2–Q3 | Near-miss ve Sınırda Geçenler")
            slide.shapes.add_picture(str(chart), PptInches(0.8), PptInches(1.2), width=PptInches(11.8))
            _add_ppt_finding(slide, _near_finding(near_table, limitations))

        slide = _new_content_slide(prs, "Q4 | ELL–ELT ve Diğer Lisans Programları")
        _add_ppt_table(slide, threshold_table, [
            ("analysis_exam_id", "Oturum"), ("threshold_group", "Grup"),
            ("N", "N"), ("PASS", "PASS"), ("FAIL", "FAIL"), ("PASS_rate", "PASS %"),
        ], max_rows=10)
        _add_ppt_finding(slide, _threshold_group_finding(threshold_table))

        if not faculty_near.empty:
            chart = temp / "faculty_near.png"
            _create_faculty_near_chart(faculty_near, chart)
            slide = _new_content_slide(prs, "Q5 | Near-miss’in Fakültelere Göre Dağılımı")
            slide.shapes.add_picture(str(chart), PptInches(0.75), PptInches(1.15), width=PptInches(11.9))
            _add_ppt_finding(slide, _faculty_near_finding(faculty_near))

        slide = _new_content_slide(prs, "Q8–Q10 | Eşik Çevresinde Beceri Profilleri")
        _add_ppt_table(slide, skill_table, [
            ("analysis_exam_id", "Oturum"), ("threshold_side", "Konum"), ("N", "N"),
            ("receptive_pct", "Rec."), ("productive_pct", "Prod."),
            ("writing_pct", "Wr."), ("speaking_pct", "Sp."),
            ("productive_minus_receptive", "P–R"),
        ], max_rows=10)
        _add_ppt_finding(slide, _skill_finding(skill_table))

        if not exam_type.empty:
            chart = temp / "exam_type.png"
            _create_exam_type_chart(exam_type, chart)
            slide = _new_content_slide(prs, "Q11 | Giriş ve Dönem İçi EPE Profilleri")
            slide.shapes.add_picture(str(chart), PptInches(0.9), PptInches(1.25), width=PptInches(11.4))
            _add_ppt_finding(slide, _exam_type_finding(exam_type))

        if not concentration.empty:
            chart = temp / "concentration.png"
            _create_concentration_chart(concentration, chart)
            slide = _new_content_slide(prs, "Q12 | Eşik Çevresindeki Yoğunluk")
            slide.shapes.add_picture(str(chart), PptInches(0.8), PptInches(1.2), width=PptInches(11.8))
            _add_ppt_finding(slide, _concentration_finding(concentration))

    slide = _new_content_slide(prs, "Dışlanan Kayıtların Özeti")
    _add_ppt_table(slide, exclusions, [
        ("analysis_exam_id", "Oturum"), ("exclusion_reason", "Neden"),
        ("N", "N"), ("PASS", "PASS"), ("FAIL", "FAIL"), ("ABSENT", "Absent"),
    ], max_rows=10)

    slide = _new_content_slide(prs, "Veri Kalitesi: Açık Kontrol Listesi")
    _add_ppt_table(slide, quality_issues, [
        ("severity", "Düzey"), ("area", "Alan"), ("file_or_exam", "Dosya/Oturum"),
        ("issue", "Sorun"), ("impact", "Etkisi"),
    ], max_rows=8)

    slide = _new_content_slide(prs, "Karar İçin Öne Çıkan Alanlar")
    textbox = slide.shapes.add_textbox(PptInches(0.9), PptInches(1.35), PptInches(11.6), PptInches(5.5))
    tf = textbox.text_frame
    prompts = [
        "Eylül near-miss sıfırları borderline süreç nedeniyle gerçek sınır yoğunluğunu yansıtıyor mu?",
        "Fakültesi hâlâ Unknown kalan kayıtlar karar öncesinde tamamlanmalı mı?",
        "Giriş ve dönem içi EPE profillerindeki fark hangi öğrenci bileşimi değişkenleriyle açıklanabilir?",
        "Productive–Receptive ve Writing–Speaking farkları hangi destek alanlarına işaret ediyor?",
        "Hangi küçük hücreler ek veri olmadan yorumlanmamalı?",
    ]
    for idx, item in enumerate(prompts):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = PptPt(22)
        p.space_after = PptPt(14)

    prs.save(path)
